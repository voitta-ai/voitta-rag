"""Parser for Microsoft PowerPoint documents (.pptx)."""

from pathlib import Path

from .base import BaseParser, ParserResult


class PptxParser(BaseParser):
    """Parser for .pptx files using python-pptx."""

    extensions = [".pptx"]

    def parse(self, file_path: Path) -> ParserResult:
        """Parse a PowerPoint presentation to Markdown."""
        try:
            from pptx import Presentation
            from pptx.exc import PackageNotFoundError
        except ImportError:
            return ParserResult.failure("python-pptx not installed. Run: pip install python-pptx")

        try:
            prs = Presentation(file_path)
        except PackageNotFoundError:
            return ParserResult.failure(f"Invalid or corrupted PPTX file: {file_path}")
        except Exception as e:
            return ParserResult.failure(f"Failed to open PPTX file: {e}")

        lines = []
        metadata = {}

        # Extract core properties
        try:
            core_props = prs.core_properties
            if core_props.title:
                metadata["title"] = core_props.title
            if core_props.author:
                metadata["author"] = core_props.author
        except Exception:
            pass

        # Add presentation title
        if metadata.get("title"):
            lines.append(f"# {metadata['title']}")
            lines.append("")

        # Process each slide
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_lines = []
            slide_title = None

            # Get shapes sorted by position (top to bottom, left to right)
            shapes = list(slide.shapes)
            shapes.sort(key=lambda s: (s.top or 0, s.left or 0))

            for shape in shapes:
                # Check if it's a title placeholder
                if shape.is_placeholder and hasattr(shape, "placeholder_format"):
                    if shape.placeholder_format.type == 1:  # Title
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_title = shape.text.strip()
                            continue

                # Extract text from shape
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # Handle bullet points (paragraphs)
                    if hasattr(shape, "text_frame"):
                        for para in shape.text_frame.paragraphs:
                            para_text = para.text.strip()
                            if para_text:
                                level = para.level if para.level else 0
                                indent = "  " * level
                                slide_lines.append(f"{indent}- {para_text}")
                    else:
                        slide_lines.append(text)

                # Handle tables
                if shape.has_table:
                    table_md = self._table_to_markdown(shape.table)
                    if table_md:
                        slide_lines.append("")
                        slide_lines.append(table_md)

            # Add slide header
            if slide_title:
                lines.append(f"## Slide {slide_num}: {slide_title}")
            else:
                lines.append(f"## Slide {slide_num}")
            lines.append("")

            # Add slide content
            if slide_lines:
                lines.extend(slide_lines)
            else:
                lines.append("*(No text content)*")

            lines.append("")
            lines.append("---")
            lines.append("")

        content = "\n".join(lines).strip()
        return ParserResult(content=content, metadata=metadata)

    def _table_to_markdown(self, table) -> str:
        """Convert a PowerPoint table to Markdown format."""
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace("\n", " ")
                cells.append(cell_text)
            rows.append(cells)

        if not rows:
            return ""

        # Build markdown table
        lines = []

        # Header row
        lines.append("| " + " | ".join(rows[0]) + " |")

        # Separator
        lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")

        # Data rows
        for row in rows[1:]:
            while len(row) < len(rows[0]):
                row.append("")
            lines.append("| " + " | ".join(row[: len(rows[0])]) + " |")

        return "\n".join(lines)
